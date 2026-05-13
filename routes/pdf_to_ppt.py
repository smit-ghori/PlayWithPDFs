from flask import Blueprint, render_template, request, send_file
from io import BytesIO

import os
import tempfile
import zipfile

import fitz  # PyMuPDF
from pptx import Presentation

pdf_to_ppt_bp = Blueprint("pdf_to_ppt", __name__)


# SHOW PAGE
@pdf_to_ppt_bp.route("/pdf_to_ppt")
def pdf_to_ppt():
    return render_template("pdf_to_ppt.html")


# CONVERT PDF TO PPT
@pdf_to_ppt_bp.route("/pdf_to_ppt", methods=["POST"])
def convert_pdf_to_ppt():

    uploaded_files = request.files.getlist("pdfs")

    if not uploaded_files:
        return "No PDF files uploaded"

    # Temporary working directory
    with tempfile.TemporaryDirectory() as temp_dir:

        generated_ppts = []

        for uploaded_file in uploaded_files:

            try:

                # Save uploaded PDF
                pdf_path = os.path.join(temp_dir, uploaded_file.filename)

                uploaded_file.save(pdf_path)

                # Create PowerPoint presentation
                presentation = Presentation()

                # Blank slide layout
                blank_slide_layout = presentation.slide_layouts[6]

                # Open PDF safely
                with fitz.open(pdf_path) as pdf_document:

                    # Convert each PDF page into PPT slide
                    for page_number in range(len(pdf_document)):

                        page = pdf_document.load_page(page_number)

                        # High quality render
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))

                        image_path = os.path.join(temp_dir, f"{page_number}.png")

                        pix.save(image_path)

                        # Create slide
                        slide = presentation.slides.add_slide(blank_slide_layout)

                        # Add image to slide
                        slide.shapes.add_picture(
                            image_path,
                            0,
                            0,
                            width=presentation.slide_width,
                            height=presentation.slide_height,
                        )

                # Output PPTX filename
                ppt_filename = os.path.splitext(uploaded_file.filename)[0] + ".pptx"

                ppt_output_path = os.path.join(temp_dir, ppt_filename)

                # Save presentation
                presentation.save(ppt_output_path)

                generated_ppts.append(ppt_output_path)

            except Exception as e:
                return f"Error converting {uploaded_file.filename}: {str(e)}"

        # SINGLE PPT DOWNLOAD
        if len(generated_ppts) == 1:

            # Read PPT into memory
            with open(generated_ppts[0], "rb") as ppt_file:

                ppt_data = ppt_file.read()

            ppt_buffer = BytesIO(ppt_data)

            ppt_buffer.seek(0)

            return send_file(
                ppt_buffer,
                as_attachment=True,
                download_name=os.path.basename(generated_ppts[0]),
                mimetype="application/octet-stream",
            )

        # MULTIPLE PPTS -> ZIP DOWNLOAD
        else:

            # Create ZIP in memory
            zip_buffer = BytesIO()

            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:

                for ppt_path in generated_ppts:

                    with open(ppt_path, "rb") as ppt_file:

                        zip_file.writestr(os.path.basename(ppt_path), ppt_file.read())

            zip_buffer.seek(0)

            return send_file(
                zip_buffer,
                as_attachment=True,
                download_name="converted_ppts.zip",
                mimetype="application/zip",
            )
